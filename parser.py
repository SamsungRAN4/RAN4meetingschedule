import pandas as pd
from pptx import Presentation
import re
import json

# ==========================================
# 1. 配置文件路径
# ==========================================
CSV_FILES = [
    "List_Topics_moderators-RAN4_119_v03_gene_r3.csv",
    "List_Topics_moderators-RAN4_119_v03_shan_r3.csv",
    "List_Topics_moderators-RAN4_119_v03_yang_r3.csv"
]
PPTX_FILE = "RAN4_119_meeting_schedule_v06_clean.pptx"
OUTPUT_FILE = "dashboard_data.js"

def parse_csv_delegates(csv_files):
    delegate_map = {}
    for file in csv_files:
        try:
            df = pd.read_csv(file)
            resp_col = [c for c in df.columns if 'Samsung Major Responsibility' in str(c)]
            mon_col = [c for c in df.columns if 'Work' in str(c) or 'Monitoring' in str(c)]
            index_col = df.columns[0]
            if not resp_col: continue
            resp_col_name = resp_col[0]
            mon_col_name = mon_col[0] if mon_col else None
            for _, row in df.iterrows():
                topic_idx = str(row[index_col]).strip()
                if topic_idx == 'nan' or not topic_idx: continue
                resps_raw = str(row[resp_col_name]).strip()
                owners = [] if resps_raw in ['nan', 'N/A', ''] else [o.strip() for o in resps_raw.split('/') if o.strip()]
                if mon_col_name and str(row[mon_col_name]).strip().lower() == 'monitoring':
                    owners = [f"{o} (monitoring)" for o in owners]
                delegate_map[topic_idx] = owners
                base_idx = topic_idx.split('-')[0]
                if base_idx not in delegate_map: delegate_map[base_idx] = []
                delegate_map[base_idx].extend(owners); delegate_map[base_idx] = list(set(delegate_map[base_idx]))
        except Exception as e: print(f"Error parsing {file}: {e}")
    return delegate_map

def parse_pptx_schedule(pptx_path):
    prs = Presentation(pptx_path)
    schedule_data = {}
    days = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"]
    day_idx = 0
    BREAK_KEYWORDS = ['lunch', 'break', '茶歇', 'dinner', 'waqfa', 'kafè']

    for slide in prs.slides:
        if day_idx >= len(days): break
        current_day = days[day_idx]; day_schedule = []
        schedule_table = None; max_rows = 0
        for shape in slide.shapes:
            if shape.has_table and len(shape.table.rows) > max_rows:
                max_rows = len(shape.table.rows); schedule_table = shape.table
        if not schedule_table: continue
        
        table = schedule_table
        room_names = ["Yang", "Shan", "Gene", "Ad-hoc"]
        
        raw_rows = []
        for row_idx in range(len(table.rows)):
            row = table.rows[row_idx]
            time_text = row.cells[0].text_frame.text.strip().replace('\n', '')
            if not time_text or not any(char.isdigit() for char in time_text): continue
            cell_texts = [cell.text_frame.text.strip() for cell in row.cells]
            while len(cell_texts) < 5: cell_texts.append("")
            raw_rows.append({"time": time_text, "cells": cell_texts})

        # Pass 2: Coffee break AH 逻辑移动 (只向下寻找)
        for i, r in enumerate(raw_rows):
            for col_idx in range(1, 5):
                text = r['cells'][col_idx]
                if not text: continue
                match = re.search(r'(?i)(coffee\s*break\s*(?:ah|ad-?hoc).*)', text, flags=re.DOTALL)
                if match:
                    ah_text = match.group(1).strip()
                    original_text = text[:match.start()].strip()
                    moved = False
                    if i < len(raw_rows) - 1:
                        next_row_text = " ".join(raw_rows[i+1]['cells']).lower()
                        if ('coffee' in next_row_text or '茶歇' in next_row_text) and 'lunch' not in next_row_text:
                            raw_rows[i+1]['cells'][col_idx] = (raw_rows[i+1]['cells'][col_idx] + "\n" + ah_text).strip()
                            moved = True
                    if moved: r['cells'][col_idx] = original_text

        # Pass 3: 提取原始休息词并填充
        last_topics_memory = {name: [] for name in room_names}
        for r in raw_rows:
            time_text = r['time']; cell_texts = r['cells']
            
            # 【核心优化】：提取该行真正使用的休息词
            current_row_break_word = "Coffee Break"
            for txt in cell_texts[1:]:
                # 寻找包含关键字但没有编号的格，作为本行的“模板词”
                if any(kw in txt.lower() for kw in BREAK_KEYWORDS) and not re.search(r'\[[0-9]+', txt):
                    current_row_break_word = txt.strip(); break
            
            has_any_topic = any(re.search(r'\[[0-9]+(?:-[A-Z])?\]', txt) for txt in cell_texts[1:])
            row_full_text = " ".join(cell_texts).lower()
            is_break_keyword = any(kw in row_full_text for kw in BREAK_KEYWORDS)
            
            if is_break_keyword and not has_any_topic:
                day_schedule.append({"time": time_text, "isBreak": True, "text": current_row_break_word.replace('\n', '<br>')})
                for nm in room_names: last_topics_memory[nm] = []
            else:
                rooms_data = []; is_split_break = is_break_keyword and has_any_topic
                for col_idx in range(1, 5):
                    raw_text = cell_texts[col_idx]; room_name = room_names[col_idx - 1]
                    found_topics = re.findall(r'\[([0-9]+(?:-[A-Z])?)\]', raw_text)
                    if is_split_break:
                        # 【核心优化】：分裂茶歇行，没任务的格子填充本行的原始休息词
                        if not found_topics: raw_text = current_row_break_word
                        last_topics_memory[room_name] = []
                    else:
                        raw_lower = raw_text.lower()
                        if any(word in raw_lower for word in ['tbd', 'early return', 'return to', 'main session']) or not raw_text:
                            found_topics = []; last_topics_memory[room_name] = []
                        elif not found_topics: found_topics = last_topics_memory[room_name].copy()
                        else: last_topics_memory[room_name] = found_topics.copy()
                    rooms_data.append({"name": room_name, "text": raw_text.replace('\n', '<br>'), "topics": found_topics})
                day_schedule.append({"time": time_text, "rooms": rooms_data})
        schedule_data[current_day] = day_schedule; day_idx += 1
    return schedule_data

if __name__ == "__main__":
    delegate_map = parse_csv_delegates(CSV_FILES)
    schedule_data = parse_pptx_schedule(PPTX_FILE)
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as f:
        f.write(f"const DELEGATE_MAP = {json.dumps(delegate_map, indent=4, ensure_ascii=False)};\n\n")
        f.write(f"const scheduleData = {json.dumps(schedule_data, indent=4, ensure_ascii=False)};\n")
    print(f"Success! Data saved to {OUTPUT_FILE}.")